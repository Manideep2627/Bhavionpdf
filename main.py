"""
Bhavion PDF Converter - FastAPI Backend (Windows Compatible)
"""

import io
import os
import sys
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import List, Optional

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response, FileResponse
from fastapi.staticfiles import StaticFiles

from pypdf import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from PIL import Image

app = FastAPI(title="Bhavion PDF API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ---
BASE_DIR = Path(__file__).parent

@app.get("/")
async def serve_frontend():
    return FileResponse(BASE_DIR / "index.html")

# Serve all 20 tool pages
TOOL_PAGES = [
    "word-to-pdf", "jpg-to-pdf", "excel-to-pdf", "powerpoint-to-pdf",
    "html-to-pdf", "pdf-to-jpg", "pdf-to-word", "pdf-to-excel",
    "pdf-to-powerpoint", "merge-pdf", "split-pdf", "compress-pdf",
    "remove-pages", "extract-pages", "organize-pdf", "rotate-pdf",
    "add-watermark", "add-page-numbers", "protect-pdf", "unlock-pdf",
]

for _slug in TOOL_PAGES:
    def make_route(slug):
        @app.get(f"/{slug}")
        async def tool_page():
            return FileResponse(BASE_DIR / f"{slug}.html")
        tool_page.__name__ = f"page_{slug.replace('-','_')}"
    make_route(_slug)

# ---
def _find_libreoffice() -> str:
    # Windows paths
    win_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for p in win_paths:
        if Path(p).exists():
            return p
    # Linux/Mac
    for cmd in ["libreoffice", "soffice"]:
        if shutil.which(cmd):
            return cmd
    return None

LIBREOFFICE = _find_libreoffice()

# ---
def _tmp_dir() -> Path:
    return Path(tempfile.mkdtemp())

def _read_upload(file: UploadFile) -> bytes:
    return file.file.read()

def _stream(data: bytes, filename: str, media: str = "application/pdf") -> Response:
    return Response(
        content=data,
        media_type=media,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )

def _libreoffice_to_pdf(src: Path, out_dir: Path) -> Path:
    if not LIBREOFFICE:
        raise HTTPException(500, "LibreOffice is not installed. Please install it from https://www.libreoffice.org/download/")
    subprocess.run(
        [LIBREOFFICE, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(src)],
        check=True,
        capture_output=True,
    )
    pdfs = list(out_dir.glob("*.pdf"))
    if not pdfs:
        raise HTTPException(500, "Conversion failed - LibreOffice produced no output")
    return pdfs[0]

def _images_to_pdf(images: List[Image.Image]) -> bytes:
    buf = io.BytesIO()
    rgb_images = [img.convert("RGB") for img in images]
    rgb_images[0].save(buf, format="PDF", save_all=True, append_images=rgb_images[1:])
    return buf.getvalue()

# ===
#  CONVERT -> PDF
# ===

@app.post("/api/convert/image-to-pdf")
async def image_to_pdf(files: List[UploadFile] = File(...)):
    images = [Image.open(io.BytesIO(_read_upload(f))) for f in files]
    return _stream(_images_to_pdf(images), "converted.pdf")

@app.post("/api/convert/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    tmp = _tmp_dir()
    try:
        src = tmp / file.filename
        src.write_bytes(_read_upload(file))
        pdf = _libreoffice_to_pdf(src, tmp)
        return _stream(pdf.read_bytes(), "converted.pdf")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

@app.post("/api/convert/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    tmp = _tmp_dir()
    try:
        src = tmp / file.filename
        src.write_bytes(_read_upload(file))
        pdf = _libreoffice_to_pdf(src, tmp)
        return _stream(pdf.read_bytes(), "converted.pdf")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

@app.post("/api/convert/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    tmp = _tmp_dir()
    try:
        src = tmp / file.filename
        src.write_bytes(_read_upload(file))
        pdf = _libreoffice_to_pdf(src, tmp)
        return _stream(pdf.read_bytes(), "converted.pdf")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

@app.post("/api/convert/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    tmp = _tmp_dir()
    try:
        src = tmp / file.filename
        src.write_bytes(_read_upload(file))
        pdf = _libreoffice_to_pdf(src, tmp)
        return _stream(pdf.read_bytes(), "converted.pdf")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

# ===
#  CONVERT FROM PDF
# ===

@app.post("/api/convert/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    try:
        import fitz  # pymupdf
    except ImportError:
        raise HTTPException(500, "Please run: pip install pymupdf")
    data = _read_upload(file)
    doc = fitz.open(stream=data, filetype="pdf")
    images = []
    for page in doc:
        mat = fitz.Matrix(2.0, 2.0)  # 2x zoom = ~150dpi
        pix = page.get_pixmap(matrix=mat)
        img_buf = io.BytesIO(pix.tobytes("jpeg"))
        images.append(img_buf.getvalue())
    doc.close()
    if len(images) == 1:
        return _stream(images[0], "page-1.jpg", "image/jpeg")
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w") as zf:
        for i, img_bytes in enumerate(images, 1):
            zf.writestr(f"page-{i}.jpg", img_bytes)
    zip_buf.seek(0)
    return _stream(zip_buf.read(), "pages.zip", "application/zip")

@app.post("/api/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    import fitz
    from docx import Document
    from docx.shared import Pt
    data = _read_upload(file)
    # Try LibreOffice first
    if LIBREOFFICE:
        tmp = _tmp_dir()
        try:
            src = tmp / "input.pdf"
            src.write_bytes(data)
            result = subprocess.run(
                [LIBREOFFICE, "--headless", "--convert-to", "docx", "--outdir", str(tmp), str(src)],
                capture_output=True, timeout=60,
            )
            docx_files = list(tmp.glob("*.docx"))
            if docx_files and docx_files[0].stat().st_size > 100:
                return _stream(
                    docx_files[0].read_bytes(), "converted.docx",
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
        except Exception:
            pass
        finally:
            shutil.rmtree(tmp, ignore_errors=True)
    # Fallback: extract text with pymupdf and write to docx
    doc_out = Document()
    pdf_doc = fitz.open(stream=data, filetype="pdf")
    for page_num, page in enumerate(pdf_doc):
        if page_num > 0:
            doc_out.add_page_break()
        blocks = page.get_text("blocks")
        for block in sorted(blocks, key=lambda b: (b[1], b[0])):
            text = block[4].strip()
            if text:
                para = doc_out.add_paragraph(text)
                para.style.font.size = Pt(11)
    pdf_doc.close()
    buf = io.BytesIO()
    doc_out.save(buf)
    buf.seek(0)
    return _stream(
        buf.read(), "converted.docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )

@app.post("/api/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    import pdfplumber, openpyxl
    data = _read_upload(file)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    row_cursor = 1
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    for row in table:
                        for col_idx, cell in enumerate(row, 1):
                            ws.cell(row=row_cursor, column=col_idx, value=cell or "")
                        row_cursor += 1
                    row_cursor += 1
            else:
                for line in (page.extract_text() or "").split("\n"):
                    ws.cell(row=row_cursor, column=1, value=line)
                    row_cursor += 1
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return _stream(buf.read(), "converted.xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.post("/api/convert/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    from pptx import Presentation
    from pptx.util import Inches
    try:
        import fitz
    except ImportError:
        raise HTTPException(500, "Please run: pip install pymupdf")
    data = _read_upload(file)
    doc = fitz.open(stream=data, filetype="pdf")
    tmp = _tmp_dir()
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img_path = tmp / f"slide_{i}.jpg"
            pix.save(str(img_path))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(img_path), 0, 0, Inches(10), Inches(7.5))
        doc.close()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return _stream(buf.read(), "converted.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

# ===
#  ORGANIZE
# ===

@app.post("/api/organize/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    writer = PdfWriter()
    for f in files:
        for page in PdfReader(io.BytesIO(_read_upload(f))).pages:
            writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "merged.pdf")

@app.post("/api/organize/split")
async def split_pdf(file: UploadFile = File(...), pages: Optional[str] = Form(None)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    total = len(reader.pages)

    def parse_ranges(spec):
        groups = []
        for part in spec.split(","):
            part = part.strip()
            if "-" in part:
                a, b = part.split("-")
                groups.append(list(range(int(a)-1, int(b))))
            else:
                groups.append([int(part)-1])
        return groups

    groups = parse_ranges(pages) if pages else [[i] for i in range(total)]
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w") as zf:
        for idx, group in enumerate(groups, 1):
            writer = PdfWriter()
            for pg in group:
                if 0 <= pg < total:
                    writer.add_page(reader.pages[pg])
            part_buf = io.BytesIO()
            writer.write(part_buf)
            part_buf.seek(0)
            zf.writestr(f"split_{idx}.pdf", part_buf.read())
    zip_buf.seek(0)
    return _stream(zip_buf.read(), "split.zip", "application/zip")

@app.post("/api/organize/remove-pages")
async def remove_pages(file: UploadFile = File(...), pages: str = Form(...)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    to_remove = {int(p.strip())-1 for p in pages.split(",")}
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        if i not in to_remove:
            writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "result.pdf")

@app.post("/api/organize/extract-pages")
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    def parse(spec):
        idxs = []
        for part in spec.split(","):
            part = part.strip()
            if "-" in part:
                a, b = part.split("-")
                idxs.extend(range(int(a)-1, int(b)))
            else:
                idxs.append(int(part)-1)
        return idxs
    writer = PdfWriter()
    for idx in parse(pages):
        if 0 <= idx < len(reader.pages):
            writer.add_page(reader.pages[idx])
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "extracted.pdf")

@app.post("/api/organize/reorder")
async def reorder_pdf(file: UploadFile = File(...), order: str = Form(...)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    total = len(reader.pages)
    writer = PdfWriter()
    for idx_str in order.split(","):
        idx = int(idx_str.strip())-1
        if 0 <= idx < total:
            writer.add_page(reader.pages[idx])
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "reordered.pdf")

# ===
#  OPTIMIZE
# ===

@app.post("/api/optimize/compress")
async def compress_pdf(
    file: UploadFile = File(...),
    level: str = Form("ebook"),  # screen=smallest, ebook=medium, printer=low
):
    import fitz, glob
    data = _read_upload(file)
    original_size = len(data)

    level_map = {
        "screen":  {"quality": 10, "scale": 0.8},
        "ebook":   {"quality": 30, "scale": 1.0},
        "printer": {"quality": 60, "scale": 1.2},
    }
    cfg = level_map.get(level, level_map["ebook"])
    gs_setting = level  # screen/ebook/printer maps directly to gs settings

    def try_ghostscript() -> Optional[bytes]:
        tmp = _tmp_dir()
        try:
            src = tmp / "input.pdf"
            src.write_bytes(data)
            out = tmp / "output.pdf"
            gs_candidates = ["gswin64c", "gswin32c", "gs"]
            for pattern in [r"C:\Program Files\gs\*\bin\gswin64c.exe",
                            r"C:\Program Files (x86)\gs\*\bin\gswin32c.exe"]:
                matches = glob.glob(pattern)
                if matches:
                    gs_candidates = matches + gs_candidates
            for gs in gs_candidates:
                try:
                    subprocess.run(
                        [gs, "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
                         f"-dPDFSETTINGS=/{gs_setting}", "-dNOPAUSE", "-dQUIET",
                         "-dBATCH", f"-sOutputFile={out}", str(src)],
                        capture_output=True, timeout=60, check=True,
                    )
                    if out.exists() and out.stat().st_size > 100:
                        return out.read_bytes()
                except (FileNotFoundError, subprocess.CalledProcessError):
                    continue
            return None
        finally:
            shutil.rmtree(tmp, ignore_errors=True)

    def try_pymupdf() -> bytes:
        quality = cfg["quality"]
        scale = cfg["scale"]
        pdf_in = fitz.open(stream=data, filetype="pdf")
        pdf_out = fitz.open()
        for page in pdf_in:
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("jpeg", jpg_quality=quality)
            tmp_img = fitz.open()
            tmp_page = tmp_img.new_page(width=pix.width, height=pix.height)
            tmp_page.insert_image(fitz.Rect(0, 0, pix.width, pix.height), stream=img_bytes)
            out_buf = io.BytesIO()
            tmp_img.save(out_buf)
            tmp_img.close()
            page_doc = fitz.open("pdf", out_buf.getvalue())
            pdf_out.insert_pdf(page_doc)
            page_doc.close()
        pdf_in.close()
        out = io.BytesIO()
        pdf_out.save(out, garbage=4, deflate=True, clean=True)
        pdf_out.close()
        out.seek(0)
        return out.read()

    def try_pypdf() -> bytes:
        reader = PdfReader(io.BytesIO(data))
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)
        buf = io.BytesIO()
        writer.write(buf)
        buf.seek(0)
        return buf.read()

    try:
        results = []
        gs_result = try_ghostscript()
        if gs_result:
            results.append(gs_result)
        img_result = try_pymupdf()
        results.append(img_result)
        results.append(try_pypdf())
        best = min(results, key=len)
        return _stream(best, "compressed.pdf")
    except Exception:
        return _stream(try_pypdf(), "compressed.pdf")


@app.post("/api/edit/rotate")
async def rotate_pdf(file: UploadFile = File(...), angle: int = Form(90), pages: Optional[str] = Form(None)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    writer = PdfWriter()
    to_rotate = {int(p.strip())-1 for p in pages.split(",")} if pages else set(range(len(reader.pages)))
    for i, page in enumerate(reader.pages):
        if i in to_rotate:
            page.rotate(angle)
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "rotated.pdf")

@app.post("/api/edit/watermark")
async def add_watermark(file: UploadFile = File(...), text: str = Form("CONFIDENTIAL"), opacity: float = Form(0.25)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    writer = PdfWriter()
    page_w = float(reader.pages[0].mediabox.width)
    page_h = float(reader.pages[0].mediabox.height)
    wm_buf = io.BytesIO()
    c = canvas.Canvas(wm_buf, pagesize=(page_w, page_h))
    c.setFont("Helvetica-Bold", 48)
    c.setFillColorRGB(0.5, 0.5, 0.5, alpha=opacity)
    c.saveState()
    c.translate(page_w/2, page_h/2)
    c.rotate(45)
    c.drawCentredString(0, 0, text)
    c.restoreState()
    c.save()
    wm_buf.seek(0)
    watermark_page = PdfReader(wm_buf).pages[0]
    for page in reader.pages:
        page.merge_page(watermark_page)
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "watermarked.pdf")

@app.post("/api/edit/add-page-numbers")
async def add_page_numbers(file: UploadFile = File(...), position: str = Form("bottom-center"), start: int = Form(1)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    writer = PdfWriter()
    for i, page in enumerate(reader.pages):
        page_w = float(page.mediabox.width)
        page_h = float(page.mediabox.height)
        overlay_buf = io.BytesIO()
        c = canvas.Canvas(overlay_buf, pagesize=(page_w, page_h))
        c.setFont("Helvetica", 11)
        c.setFillColorRGB(0.2, 0.2, 0.2)
        num = str(i + start)
        if position == "bottom-right":
            c.drawRightString(page_w-30, 20, num)
        elif position == "top-center":
            c.drawCentredString(page_w/2, page_h-25, num)
        else:
            c.drawCentredString(page_w/2, 20, num)
        c.save()
        overlay_buf.seek(0)
        page.merge_page(PdfReader(overlay_buf).pages[0])
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "numbered.pdf")

# ===
#  SECURITY
# ===

@app.post("/api/security/protect")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    reader = PdfReader(io.BytesIO(_read_upload(file)))
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(user_password=password, owner_password=password+"_owner")
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "protected.pdf")

@app.post("/api/security/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form("")):
    data = _read_upload(file)
    reader = PdfReader(io.BytesIO(data))
    if reader.is_encrypted:
        if not reader.decrypt(password):
            raise HTTPException(400, "Wrong password")
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    buf = io.BytesIO()
    writer.write(buf)
    buf.seek(0)
    return _stream(buf.read(), "unlocked.pdf")

# ---
@app.get("/api/health")
async def health():
    lo = "Found" if LIBREOFFICE else "NOT FOUND - please install LibreOffice from https://www.libreoffice.org"
    return {"status": "ok", "libreoffice": lo}
